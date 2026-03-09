#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script de création de présentation PowerPoint pour GesRes
Gestion des Missions Réserve - Gendarmerie Nationale
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Couleurs de la Gendarmerie
BLEU_GENDARMERIE = RGBColor(0, 51, 127)  # Bleu foncé
BLEU_CLAIR = RGBColor(41, 98, 255)
BLANC = RGBColor(255, 255, 255)
GRIS_FONCE = RGBColor(51, 51, 51)
GRIS_CLAIR = RGBColor(240, 240, 240)

def create_title_slide(prs):
    """Slide 1 : Page de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond bleu
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLEU_GENDARMERIE
    
    # Titre principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "GesRes"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = BLANC
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Gestion des Missions Réserve"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(36)
    subtitle_para.font.color.rgb = BLANC
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Gendarmerie Nationale - 2026"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(18)
    footer_para.font.color.rgb = BLANC

def create_context_slide(prs):
    """Slide 2 : Contexte et problématique"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Titre
    title = slide.shapes.title
    title.text = "📋 Contexte et Problématique"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    # Contenu
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    points = [
        ("Besoin", "Optimiser la gestion des missions pour les réservistes de la Gendarmerie"),
        ("Problème actuel", "Gestion manuelle, difficultés de coordination, manque de visibilité"),
        ("Objectif", "Digitaliser et automatiser la gestion des missions et des affectations"),
        ("Bénéfices", "Gain de temps, meilleure organisation, suivi en temps réel")
    ]
    
    for label, desc in points:
        p = text_frame.add_paragraph()
        p.text = f"• {label} : "
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.level = 1
        p2.font.size = Pt(16)
        p2.space_after = Pt(12)

def create_solution_slide(prs):
    """Slide 3 : Solution proposée"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💡 Solution : GesRes"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = "Application web 100% responsive"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLEU_GENDARMERIE
    p.space_after = Pt(18)
    
    features = [
        "📱 Accessible sur smartphone, tablette et ordinateur",
        "🌐 Déployée sur Cloudflare Pages (rapide et sécurisé)",
        "🔐 Authentification sécurisée avec JWT",
        "📊 Base de données D1 SQLite distribuée",
        "⚡ Interface moderne et intuitive",
        "🗑️ Nettoyage automatique des missions expirées"
    ]
    
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(18)
        p.space_after = Pt(10)

def create_features_public_slide(prs):
    """Slide 4 : Fonctionnalités - Interface publique"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "👥 Interface Publique (Gendarmes)"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    sections = [
        ("Navigation hiérarchique intuitive", [
            "Sélection de la compagnie",
            "Choix de la brigade",
            "Consultation des missions par mois",
            "Fil d'ariane pour navigation rapide"
        ]),
        ("Informations complètes sur les missions", [
            "Titre, description et lieu",
            "Dates de début/fin et durée",
            "Priorité et effectifs requis",
            "Gendarmes assignés (validés/en attente)"
        ]),
        ("Filtres avancés", [
            "Par priorité (haute/moyenne/normale)",
            "Par disponibilité (places libres)",
            "Recherche textuelle"
        ])
    ]
    
    for section_title, items in sections:
        p = text_frame.add_paragraph()
        p.text = section_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(8)
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(14)
            p.space_after = Pt(4)
        
        p.space_after = Pt(12)

def create_features_admin_slide(prs):
    """Slide 5 : Fonctionnalités - Interface admin"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "⚙️ Interface Administration"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    tabs = [
        ("📋 Gestion des Missions", [
            "Création/modification/suppression",
            "Organisation par mois",
            "Export PDF par mois"
        ]),
        ("🗺️ Gestion des Lieux", [
            "Compagnies et brigades",
            "Informations de contact",
            "Hiérarchie organisationnelle"
        ]),
        ("👤 Gestion des Gendarmes", [
            "Matricule, grade, spécialité",
            "Disponibilité",
            "Recherche avancée"
        ]),
        ("📊 Gestion des Assignations", [
            "Affectation gendarmes/missions",
            "Statuts (libre/en attente/validé)",
            "Suivi des effectifs"
        ])
    ]
    
    for tab_title, items in tabs:
        p = text_frame.add_paragraph()
        p.text = tab_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(6)
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(13)
            p.space_after = Pt(3)
        
        p.space_after = Pt(8)

def create_mobile_slide(prs):
    """Slide 6 : Version mobile/tablette"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📱 Version Mobile & Tablette"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = "Interface 100% responsive (v5.0)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = BLEU_GENDARMERIE
    p.space_after = Pt(18)
    
    devices = [
        ("📱 Smartphone (< 640px)", [
            "Menu hamburger avec overlay",
            "Grilles 1 colonne",
            "Zones tactiles 44px minimum (norme WCAG)",
            "Modales plein écran",
            "Navigation sticky"
        ]),
        ("📱 Tablette (640px - 1024px)", [
            "Grilles 2 colonnes",
            "Navigation classique",
            "Modales adaptées (90% largeur)"
        ]),
        ("💻 Desktop (> 1024px)", [
            "Grilles 3-4 colonnes",
            "Interface complète",
            "Toutes fonctionnalités disponibles"
        ])
    ]
    
    for device_title, features in devices:
        p = text_frame.add_paragraph()
        p.text = device_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(6)
        
        for feature in features:
            p = text_frame.add_paragraph()
            p.text = f"• {feature}"
            p.level = 1
            p.font.size = Pt(13)
            p.space_after = Pt(3)
        
        p.space_after = Pt(10)

def create_auto_cleanup_slide(prs):
    """Slide 7 : Nettoyage automatique"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🗑️ Nettoyage Automatique"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = "Suppression automatique des missions expirées"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLEU_GENDARMERIE
    p.space_after = Pt(16)
    
    points = [
        ("Fonctionnement", "Suppression automatique le 1er de chaque mois"),
        ("Règle", "Missions terminées avant le mois précédent sont supprimées"),
        ("Exemple", "Le 1er avril → toutes les missions de mars (terminées) disparaissent"),
        ("Gestion manuelle", "Interface admin : vérification et suppression manuelle possible"),
        ("Sécurité", "Assignations supprimées automatiquement (CASCADE)"),
        ("Logs", "Traçabilité complète des suppressions")
    ]
    
    for label, desc in points:
        p = text_frame.add_paragraph()
        p.text = f"• {label} : "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.level = 1
        p2.font.size = Pt(14)
        p2.space_after = Pt(10)

def create_architecture_slide(prs):
    """Slide 8 : Architecture technique"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🔧 Architecture Technique"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    stack = [
        ("Frontend", [
            "HTML5 / CSS3 (Tailwind CSS)",
            "JavaScript ES6+ (Vanilla JS)",
            "Font Awesome Icons",
            "Day.js pour les dates",
            "jsPDF pour l'export PDF"
        ]),
        ("Backend", [
            "Hono Framework (TypeScript)",
            "Cloudflare Workers",
            "API REST",
            "JWT Authentication"
        ]),
        ("Base de données", [
            "Cloudflare D1 (SQLite distribué)",
            "Migrations SQL versionnées",
            "Relations CASCADE"
        ]),
        ("Déploiement", [
            "Cloudflare Pages",
            "Edge Computing global",
            "HTTPS automatique",
            "CDN intégré"
        ])
    ]
    
    for section_title, items in stack:
        p = text_frame.add_paragraph()
        p.text = section_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(6)
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(13)
            p.space_after = Pt(2)
        
        p.space_after = Pt(10)

def create_benefits_slide(prs):
    """Slide 9 : Bénéfices"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "✅ Bénéfices"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    benefits = [
        ("⏱️ Gain de temps", "Automatisation de la gestion, réduction des tâches manuelles"),
        ("📊 Visibilité en temps réel", "Suivi instantané des missions et affectations"),
        ("📱 Accessibilité", "Consultable partout, sur tous les appareils"),
        ("🔒 Sécurité", "Authentification JWT, données chiffrées"),
        ("🚀 Performance", "Edge computing Cloudflare, latence minimale"),
        ("💰 Coûts réduits", "Infrastructure serverless, pas de serveur à maintenir"),
        ("♻️ Maintenance simplifiée", "Nettoyage automatique, mises à jour facilitées"),
        ("👥 Meilleure coordination", "Communication optimisée entre services")
    ]
    
    for emoji_title, desc in benefits:
        p = text_frame.add_paragraph()
        p.text = emoji_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(4)
        
        p2 = text_frame.add_paragraph()
        p2.text = desc
        p2.level = 1
        p2.font.size = Pt(14)
        p2.space_after = Pt(10)

def create_roadmap_slide(prs):
    """Slide 10 : Feuille de route"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🗺️ Feuille de Route"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    phases = [
        ("✅ Phase 1 - Terminée", [
            "Interface publique et admin",
            "Gestion missions/brigades/gendarmes",
            "Version mobile responsive",
            "Nettoyage automatique"
        ]),
        ("🚀 Phase 2 - Court terme (3 mois)", [
            "PWA (Progressive Web App)",
            "Notifications push",
            "Mode hors ligne",
            "Cron jobs Cloudflare"
        ]),
        ("💡 Phase 3 - Moyen terme (6 mois)", [
            "Tableau de bord analytique",
            "Rapports et statistiques",
            "Export Excel avancé",
            "Intégration email automatique"
        ]),
        ("🔮 Phase 4 - Long terme (12 mois)", [
            "API publique",
            "Mobile app native (iOS/Android)",
            "Intégration planning",
            "IA pour suggestions d'affectations"
        ])
    ]
    
    for phase_title, items in phases:
        p = text_frame.add_paragraph()
        p.text = phase_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(6)
        
        for item in items:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
            p.font.size = Pt(13)
            p.space_after = Pt(2)
        
        p.space_after = Pt(10)

def create_demo_slide(prs):
    """Slide 11 : Démonstration"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🎬 Démonstration"
    title.text_frame.paragraphs[0].font.color.rgb = BLEU_GENDARMERIE
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    body_shape = slide.shapes.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    p = text_frame.add_paragraph()
    p.text = "URLs d'accès"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLEU_GENDARMERIE
    p.space_after = Pt(20)
    
    urls = [
        ("🌐 Interface publique", "https://gesres-missions.pages.dev"),
        ("🔐 Interface admin", "https://gesres-missions.pages.dev/admin"),
        ("📱 Testable sur mobile", "Scanner le QR code ou ouvrir l'URL"),
        ("", ""),
        ("👤 Identifiants de test", ""),
        ("   • Utilisateur", "admin"),
        ("   • Mot de passe", "admin123")
    ]
    
    for label, value in urls:
        if not label:
            p = text_frame.add_paragraph()
            p.space_after = Pt(10)
            continue
            
        p = text_frame.add_paragraph()
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True if not label.startswith("   ") else False
        p.font.color.rgb = BLEU_GENDARMERIE
        p.space_after = Pt(6)
        
        if value:
            p2 = text_frame.add_paragraph()
            p2.text = value
            p2.level = 1
            p2.font.size = Pt(16)
            p2.font.color.rgb = GRIS_FONCE if not label.startswith("   ") else BLEU_CLAIR
            p2.space_after = Pt(10)

def create_conclusion_slide(prs):
    """Slide 12 : Conclusion"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond bleu
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLEU_GENDARMERIE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Merci pour votre attention"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = BLANC
    
    # Questions
    question_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    question_frame = question_box.text_frame
    question_frame.text = "Des questions ?"
    question_para = question_frame.paragraphs[0]
    question_para.alignment = PP_ALIGN.CENTER
    question_para.font.size = Pt(36)
    question_para.font.color.rgb = BLANC
    
    # Contact
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "🌐 gesres-missions.pages.dev\n📧 contact@gendarmerie.fr"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = Pt(20)
    contact_para.font.color.rgb = BLANC

def main():
    """Fonction principale"""
    print("🎨 Création de la présentation PowerPoint GesRes...")
    
    # Créer la présentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Créer les slides
    print("📄 Slide 1 : Page de titre")
    create_title_slide(prs)
    
    print("📄 Slide 2 : Contexte et problématique")
    create_context_slide(prs)
    
    print("📄 Slide 3 : Solution proposée")
    create_solution_slide(prs)
    
    print("📄 Slide 4 : Interface publique")
    create_features_public_slide(prs)
    
    print("📄 Slide 5 : Interface admin")
    create_features_admin_slide(prs)
    
    print("📄 Slide 6 : Version mobile")
    create_mobile_slide(prs)
    
    print("📄 Slide 7 : Nettoyage automatique")
    create_auto_cleanup_slide(prs)
    
    print("📄 Slide 8 : Architecture technique")
    create_architecture_slide(prs)
    
    print("📄 Slide 9 : Bénéfices")
    create_benefits_slide(prs)
    
    print("📄 Slide 10 : Feuille de route")
    create_roadmap_slide(prs)
    
    print("📄 Slide 11 : Démonstration")
    create_demo_slide(prs)
    
    print("📄 Slide 12 : Conclusion")
    create_conclusion_slide(prs)
    
    # Sauvegarder
    output_file = '/home/user/webapp/GesRes_Presentation.pptx'
    prs.save(output_file)
    print(f"✅ Présentation créée : {output_file}")
    print(f"📊 Nombre de slides : {len(prs.slides)}")

if __name__ == "__main__":
    main()
